from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
from .watsonx_utils import call_llm
from pathlib import Path
from uuid import uuid4
from typing import Dict, Any, List

class PowerPointGenerator:
    root_path: str = "presentations"
    main_prompt: str = "<s> [INST] قم بإنشاء عرض تقديمي لدرس {} وهذه تفاصيل الدرس: {} يتضمن العناوين الرئيسية وشرح مفصل لكل عنوان. يجب أن يحتوي كل قسم على عنوان وشرح تفصيلي باللغة العربية [/INST]"
    
    def __init__(self, section: str, description: str):
        self.section = section
        self.description = description
        # Create root directory if it doesn't exist
        Path(self.root_path).mkdir(parents=True, exist_ok=True)

    def parse_outline(self, resp_text: str) -> Dict[str, Any]:
        """Parse the LLM response into a structured outline."""
        lines = resp_text.splitlines()
        presentation_content = {}
        current_title = None
        current_points = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Check if it's a main title (Arabic title followed by a colon)
            title_match = re.match(r'^عنوان الشريحة \d+:\s*(.*)', line)
            if title_match:
                # Save previous section if exists
                if current_title:
                    presentation_content[current_title] = [c.strip(' -') for c in current_points]

                # Extract the title and reset points
                current_title = title_match.group(1).strip()
                current_points = []
            else:
                # Check if it's a bullet point (starts with - or •)
                point_match = re.match(r'^[-•]\s*(.*)', line)
                if point_match:
                    current_points.append(point_match.group(1))
                # Handle sub-points or lines that are part of the main points
                elif line.strip():
                    current_points.append(line)

        # Add the last section
        if current_title:
            presentation_content[current_title] = [c.strip(' -') for c in current_points]

        return presentation_content

    def create_presentation(self, content: Dict[str, List[str]], filename: str = "presentation"):
        """Create a PowerPoint presentation with the parsed content."""
        prs = Presentation()
        
        # Set slide width and height (16:9 aspect ratio)
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        # Create title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        title.text = self.section
        
        # Configure title formatting
        title_tf = title.text_frame
        title_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT  # Set to right alignment
        title_tf.paragraphs[0].font.size = Pt(44)
        title_tf.paragraphs[0].font.name = "Arial"
        title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Content slides
        content_slide_layout = prs.slide_layouts[1]  # Layout with title and content
        
        for slide_title, points in content.items():
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_title
            title_tf = title.text_frame
            title_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT  # Set to right alignment
            title_tf.paragraphs[0].font.size = Pt(40)
            title_tf.paragraphs[0].font.name = "Arial"
            
            # Add content
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()  # Clear existing content
            
            for point in points:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(24)
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.RIGHT  # Set each bullet point to right alignment
                p.level = 0
        
        # Save the presentation
        output_path = Path(self.root_path) / f"{filename}.pptx"
        prs.save(str(output_path))
        return output_path

    def __call__(self) -> Path:
        """Generate the presentation and return the path to the file."""
        # Generate the prompt and get LLM response
        prompt = self.main_prompt.format(self.section, self.description)
        print(f"Generating presentation for: {self.section}")
        
        # Call LLM (assuming similar function exists)
        text = call_llm(prompt, temperature = 0)
        
        # Parse the content
        presentation_content = self.parse_outline(text)
        print(f"Parsed {len(presentation_content)} sections")
        
        # Generate unique filename
        filename = str(uuid4())
        
        # Create and save the presentation
        output_path = self.create_presentation(presentation_content, filename)
        print(f"Presentation saved as: {output_path}")
        
        return output_path